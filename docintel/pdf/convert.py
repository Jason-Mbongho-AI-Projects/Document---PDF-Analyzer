"""
Conversion engine.

Two honesty rules govern this module.

Fidelity is stated, not implied. A PDF is a description of marks on a page; it
does not carry paragraphs, headings or table semantics. Anything that turns it
into DOCX or XLSX is reconstructing structure that was never recorded, and the
result will not round-trip. Every converter therefore declares its fidelity
("exact", "structural", "text-only") and that value is returned with the
output, so nobody mistakes a text dump for a faithful Word document.

Unavailable is reported, never faked. Office formats (DOCX/XLSX/PPTX) can only
be rendered to PDF by a real office engine. If LibreOffice is not installed,
those conversions are reported as unavailable with the reason, rather than
emitting a plausible-looking file that silently lost the layout.

Licences of the libraries used here — pypdf (BSD), pdfplumber (MIT),
python-docx (MIT), openpyxl (MIT), reportlab (BSD), Pillow (MIT-CMU),
pypdfium2 (BSD/Apache) — are all fine for commercial distribution.
"""
import csv
import html as html_module
import io
import json
import os
import re
import shutil
import subprocess
import zipfile
from dataclasses import dataclass, field
from typing import Callable, Dict, List, Optional, Sequence, Tuple

from docintel.pdf.engine import PDFEngineError, PasswordRequired
from docintel.pdf.render import render_page
from docintel.pdf.text import extract

# ------------------------------------------------------------- descriptors

FIDELITY = {
    "exact": "Content is reproduced exactly.",
    "structural": "Text and detected structure are preserved; visual layout is not.",
    "text-only": "Text is preserved. Layout, fonts and images are not.",
    "raster": "Pages are reproduced as images; text is no longer selectable.",
}


@dataclass
class ConversionResult:
    data: bytes
    filename: str
    media_type: str
    target: str
    fidelity: str
    note: str
    pages: int = 0
    warnings: List[str] = field(default_factory=list)

    @property
    def fidelity_note(self) -> str:
        return FIDELITY.get(self.fidelity, "")


@dataclass
class Capability:
    target: str
    label: str
    media_type: str
    extension: str
    fidelity: str
    available: bool = True
    reason: Optional[str] = None


# ------------------------------------------------------------ office tools

def _office_binary() -> Optional[str]:
    """Locate a LibreOffice binary, if one is installed."""
    for name in ("soffice", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found

    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if os.path.exists(candidate):
            return candidate
    return None


OFFICE_UNAVAILABLE = (
    "Converting Word, Excel and PowerPoint files to PDF requires a LibreOffice "
    "installation, which is not present on this server. Install LibreOffice and "
    "restart to enable it."
)


# ------------------------------------------------------- PDF -> text forms

def _pages_text(data: bytes) -> List[Tuple[int, str]]:
    return [(page.page, page.text) for page in extract(data)]


def to_txt(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    pages = _pages_text(data)
    body = "\n\n".join(
        f"--- Page {number} ---\n{text}".rstrip() for number, text in pages
    )
    return body.encode("utf-8"), len(pages), []


def to_markdown(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    """Best-effort Markdown.

    Headings are inferred from short, title-like lines. This is a heuristic:
    PDFs do not record heading levels, so treat the structure as a suggestion.
    """
    pages = _pages_text(data)
    lines_out: List[str] = []

    for number, text in pages:
        lines_out.append(f"\n<!-- page {number} -->\n")
        for raw in text.splitlines():
            line = raw.strip()
            if not line:
                continue
            words = line.split()
            looks_like_heading = (
                len(words) <= 9
                and not line.endswith((".", ",", ";", ":"))
                and (line.isupper() or line.istitle())
            )
            lines_out.append(f"## {line}" if looks_like_heading else line)
            lines_out.append("")

    return "\n".join(lines_out).strip().encode("utf-8"), len(pages), [
        "Headings are inferred heuristically; PDFs do not record heading levels."
    ]


def to_html(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    pages = _pages_text(data)
    blocks = []
    for number, text in pages:
        paragraphs = "".join(
            f"<p>{html_module.escape(line.strip())}</p>"
            for line in text.splitlines() if line.strip()
        )
        blocks.append(
            f'<section class="page" id="page-{number}">'
            f'<h2>Page {number}</h2>{paragraphs}</section>'
        )

    document = (
        "<!doctype html><html><head><meta charset=\"utf-8\">"
        "<title>Converted document</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:52rem;margin:2rem auto;"
        "line-height:1.6;padding:0 1rem}.page{margin-bottom:3rem}"
        "h2{font-size:.8rem;text-transform:uppercase;letter-spacing:.1em;color:#64748b}"
        "</style></head><body>" + "".join(blocks) + "</body></html>"
    )
    return document.encode("utf-8"), len(pages), []


def to_json(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    """Structured export: pages, words with geometry, and detected tables."""
    pages = extract(data)
    tables = _detect_tables(data)
    by_page: Dict[int, list] = {}
    for table in tables:
        by_page.setdefault(table["page"], []).append(table["rows"])

    payload = {
        "pages": [
            {
                **page.as_dict(),
                "tables": by_page.get(page.page, []),
            }
            for page in pages
        ],
        "page_count": len(pages),
    }
    return json.dumps(payload, indent=2).encode("utf-8"), len(pages), []



# XML — and therefore DOCX, XLSX and PPTX — accepts only a subset of Unicode.
# Control characters below 0x20, the surrogate range, and the noncharacters
# 0xFFFE/0xFFFF are all forbidden, and the libraries refuse the whole string
# rather than the offending character.
#
# Extracted PDF text contains them more often than one would like: a glyph
# that fails to map to Unicode frequently comes back as 0xFFFE. One such
# character in a fifty-page document used to fail the entire export with
# "All strings must be XML compatible", which tells the user nothing about
# what to do.
def _xml_ok(ch: str) -> bool:
    codepoint = ord(ch)
    return (codepoint in (0x09, 0x0A, 0x0D)
            or 0x20 <= codepoint <= 0xD7FF
            or 0xE000 <= codepoint <= 0xFFFD
            or 0x10000 <= codepoint <= 0x10FFFF)


def xml_safe(value) -> str:
    """Text with the characters XML cannot carry removed.

    They are dropped rather than replaced: an unmappable glyph carries no
    meaning to preserve, and substituting a space or a marker would invent
    content that was never in the document.
    """
    text = "" if value is None else str(value)
    return "".join(ch for ch in text if _xml_ok(ch))


# ---------------------------------------------------------------- tables

def _detect_tables(data: bytes) -> List[dict]:
    """Find tables with pdfplumber's ruling/edge detection."""
    import pdfplumber

    found: List[dict] = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                try:
                    tables = page.extract_tables() or []
                except Exception:
                    continue
                for order, rows in enumerate(tables, start=1):
                    cleaned = [
                        [("" if cell is None else str(cell).strip()) for cell in row]
                        for row in rows
                        if any(cell not in (None, "") for cell in row)
                    ]
                    if len(cleaned) >= 2:
                        found.append({"page": index, "index": order, "rows": cleaned})
    except Exception as exc:
        raise PDFEngineError(f"Tables could not be read: {exc}") from exc
    return found


def to_csv(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    tables = _detect_tables(data)
    if not tables:
        raise PDFEngineError(
            "No tables were detected in this document, so there is nothing to "
            "export as CSV. Try exporting as text or JSON instead."
        )

    buffer = io.StringIO()
    writer = csv.writer(buffer)
    for table in tables:
        writer.writerow([f"# Page {table['page']} — table {table['index']}"])
        writer.writerows(table["rows"])
        writer.writerow([])

    warnings = []
    if len(tables) > 1:
        warnings.append(
            f"{len(tables)} tables were found and written one after another "
            "into a single CSV. Export as XLSX to keep them on separate sheets."
        )
    return buffer.getvalue().encode("utf-8"), len(tables), warnings


def to_xlsx(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    tables = _detect_tables(data)
    if not tables:
        raise PDFEngineError(
            "No tables were detected in this document, so there is nothing to "
            "export as a spreadsheet."
        )

    workbook = Workbook()
    workbook.remove(workbook.active)

    for table in tables:
        sheet = workbook.create_sheet(f"p{table['page']}_t{table['index']}"[:31])
        for row in table["rows"]:
            # Numeric-looking cells are written as numbers so the sheet is
            # usable for arithmetic rather than being all text.
            sheet.append([_coerce(cell) for cell in row])
        for cell in sheet[1]:
            cell.font = Font(bold=True)

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue(), len(tables), [
        "Table boundaries are detected from ruling lines and text alignment; "
        "check the result against the original."
    ]


NUMERIC = re.compile(r"^-?[\d,]+(?:\.\d+)?$")


def _coerce(value: str):
    # Cells reach openpyxl through here, and it rejects XML-invalid text the
    # same way python-docx does.
    text = xml_safe(value).strip()
    if NUMERIC.match(text):
        try:
            return float(text.replace(",", "")) if "." in text else int(text.replace(",", ""))
        except ValueError:
            return text
    return text


# ----------------------------------------------------------------- docx

def to_docx(data: bytes, **_) -> Tuple[bytes, int, List[str]]:
    from docx import Document
    from docx.shared import Pt

    pages = _pages_text(data)
    tables_by_page: Dict[int, list] = {}
    try:
        for table in _detect_tables(data):
            tables_by_page.setdefault(table["page"], []).append(table["rows"])
    except PDFEngineError:
        tables_by_page = {}

    document = Document()
    style = document.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    for index, (number, text) in enumerate(pages):
        if index:
            document.add_page_break()

        for line in text.splitlines():
            stripped = line.strip()
            if stripped:
                document.add_paragraph(xml_safe(stripped))

        for rows in tables_by_page.get(number, []):
            if not rows:
                continue
            table = document.add_table(rows=len(rows), cols=max(len(r) for r in rows))
            table.style = "Table Grid"
            for row_index, row in enumerate(rows):
                for column_index, cell in enumerate(row):
                    table.cell(row_index, column_index).text = xml_safe(cell)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue(), len(pages), [
        "Word output contains the text and any detected tables. Original "
        "layout, fonts, images and positioning are not reproduced."
    ]


# ---------------------------------------------------------------- images

def to_images(data: bytes, fmt: str = "png", scale: float = 2.0,
              **_) -> Tuple[bytes, int, List[str]]:
    """Render every page and return them zipped."""
    pages = extract(data)
    if not pages:
        raise PDFEngineError("This document has no pages to render.")

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for page in pages:
            image = render_page(data, page.page, scale=scale, fmt=fmt)
            archive.writestr(f"page-{page.page:04d}.{fmt}", image.data)

    return buffer.getvalue(), len(pages), []


# -------------------------------------------------------------- to PDF

def text_to_pdf(text: str, *, title: str = "", page_size: str = "letter") -> bytes:
    from reportlab.lib.pagesizes import A4, LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    sizes = {"letter": LETTER, "a4": A4}
    if page_size not in sizes:
        raise PDFEngineError(f"Unknown page size '{page_size}'.")

    buffer = io.BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=sizes[page_size],
                                 title=title or "Document")
    styles = getSampleStyleSheet()

    flow = []
    if title:
        flow += [Paragraph(html_module.escape(title), styles["Title"]), Spacer(1, 12)]

    for block in text.split("\n\n"):
        block = block.strip()
        if not block:
            continue
        # Escape first: user content must never be interpreted as reportlab
        # markup.
        safe = html_module.escape(block).replace("\n", "<br/>")
        flow += [Paragraph(safe, styles["BodyText"]), Spacer(1, 8)]

    if not flow:
        flow = [Paragraph("(empty document)", styles["BodyText"])]

    document.build(flow)
    return buffer.getvalue()


def csv_to_pdf(raw: bytes, *, title: str = "") -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape, LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    try:
        rows = list(csv.reader(io.StringIO(raw.decode("utf-8-sig", errors="replace"))))
    except Exception as exc:
        raise PDFEngineError(f"The CSV could not be read: {exc}") from exc

    rows = [r for r in rows if any(cell.strip() for cell in r)]
    if not rows:
        raise PDFEngineError("The CSV is empty.")

    buffer = io.BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=landscape(LETTER),
                                 title=title or "Spreadsheet")
    styles = getSampleStyleSheet()

    flow = []
    if title:
        flow += [Paragraph(html_module.escape(title), styles["Title"]), Spacer(1, 10)]

    table = Table(rows, repeatRows=1)
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1D4ED8")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#CFDAEA")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
    ]))
    flow.append(table)

    document.build(flow)
    return buffer.getvalue()


def images_to_pdf(images: Sequence[bytes], *, page_size: str = "fit") -> bytes:
    from PIL import Image
    from reportlab.lib.pagesizes import A4, LETTER
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas as pdf_canvas

    if not images:
        raise PDFEngineError("No images were supplied.")

    buffer = io.BytesIO()
    pdf = pdf_canvas.Canvas(buffer)

    for raw in images:
        try:
            image = Image.open(io.BytesIO(raw))
            image.load()
        except Exception as exc:
            raise PDFEngineError(f"An image could not be read: {exc}") from exc

        if image.mode not in ("RGB", "L"):
            image = image.convert("RGB")

        if page_size == "fit":
            width, height = image.size
            pdf.setPageSize((width, height))
            pdf.drawImage(ImageReader(image), 0, 0, width, height)
        else:
            page = {"letter": LETTER, "a4": A4}.get(page_size)
            if page is None:
                raise PDFEngineError(f"Unknown page size '{page_size}'.")
            pdf.setPageSize(page)
            scale = min(page[0] / image.width, page[1] / image.height) * 0.94
            width, height = image.width * scale, image.height * scale
            pdf.drawImage(ImageReader(image),
                          (page[0] - width) / 2, (page[1] - height) / 2,
                          width, height)
        pdf.showPage()

    pdf.save()
    return buffer.getvalue()


def _office_environment() -> dict:
    """A child environment LibreOffice's own Python will not trip over.

    LibreOffice ships its own interpreter. Launched from inside a virtualenv it
    inherits PYTHONHOME and PYTHONPATH pointing at ours, and reports "Could not
    find platform independent libraries" — sometimes only as a warning, but it
    is a real source of first-run failures.
    """
    env = dict(os.environ)
    for name in ("PYTHONHOME", "PYTHONPATH", "PYTHONSTARTUP", "PYTHONEXECUTABLE"):
        env.pop(name, None)
    return env


def office_to_pdf(raw: bytes, extension: str) -> bytes:
    """Convert an office document via LibreOffice, if available."""
    binary = _office_binary()
    if binary is None:
        raise PDFEngineError(OFFICE_UNAVAILABLE)

    import tempfile
    from pathlib import Path

    with tempfile.TemporaryDirectory() as workdir:
        source = Path(workdir) / f"input.{extension}"
        source.write_bytes(raw)
        produced = Path(workdir) / "input.pdf"

        # A private profile per conversion. Two things depend on it: the very
        # first run on a fresh install has to create a profile, and it fails
        # rather than doing so silently; and concurrent conversions sharing the
        # default profile contend for its lock, which matters as soon as the
        # worker runs jobs in parallel.
        profile = Path(workdir) / "profile"
        profile.mkdir()
        command = [
            binary,
            f"-env:UserInstallation={profile.as_uri()}",
            "--headless", "--norestore", "--invisible", "--nologo",
            "--nolockcheck", "--nodefault",
            "--convert-to", "pdf", "--outdir", workdir, str(source),
        ]

        last_error = ""
        # Two attempts: profile creation can consume the first invocation, and
        # a single retry is cheaper than telling the user their document is
        # unconvertible when it is not.
        for attempt in (1, 2):
            try:
                completed = subprocess.run(
                    command, capture_output=True, timeout=180,
                    env=_office_environment(),
                )
            except subprocess.TimeoutExpired as exc:
                raise PDFEngineError("The conversion timed out.") from exc

            if produced.exists() and produced.stat().st_size > 0:
                return produced.read_bytes()

            last_error = (
                completed.stderr.decode(errors="replace").strip()
                or completed.stdout.decode(errors="replace").strip()
                or f"exit code {completed.returncode}"
            )

        raise PDFEngineError(
            f"The office converter could not read this .{extension} file "
            f"after {attempt} attempts: {last_error[:200]}"
        )


def to_pptx(data: bytes, scale: float = 2.0, **_) -> Tuple[bytes, int, List[str]]:
    """One slide per page, each carrying the page as a picture.

    Rendered rather than reconstructed. A PDF records glyph positions, not
    the shapes, text boxes and layout a slide is made of, so anything claiming
    to rebuild an editable deck is guessing. A faithful picture per slide is
    honest and is what the fidelity note says.

    The speaker notes carry the page's text, so the deck is still searchable
    and the words are not lost to an image.
    """
    from pptx import Presentation
    from pptx.util import Emu

    pages = extract(data)
    if not pages:
        raise PDFEngineError("This document has no pages.")

    deck = Presentation()
    blank = deck.slide_layouts[6]          # completely empty layout
    warnings: List[str] = []

    for index, page in enumerate(pages, start=1):
        # Match the slide to the page so nothing is cropped or letterboxed.
        deck.slide_width = Emu(int(page.width * 12700))
        deck.slide_height = Emu(int(page.height * 12700))

        image = render_page(data, index, scale=scale, fmt="png").data
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(image), 0, 0,
            width=deck.slide_width, height=deck.slide_height,
        )

        text = (page.text or "").strip()
        if text:
            slide.notes_slide.notes_text_frame.text = xml_safe(text)

    if len({(p.width, p.height) for p in pages}) > 1:
        warnings.append(
            "The document mixes page sizes; each slide matches its own page, "
            "so the deck is not a single consistent shape."
        )

    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue(), len(pages), warnings


# --------------------------------------------------------------- registry

Converter = Callable[..., Tuple[bytes, int, List[str]]]

_FROM_PDF: Dict[str, Tuple[Converter, Capability]] = {
    "txt": (to_txt, Capability("txt", "Plain text", "text/plain", "txt", "text-only")),
    "markdown": (to_markdown, Capability("markdown", "Markdown", "text/markdown", "md", "text-only")),
    "html": (to_html, Capability("html", "HTML", "text/html", "html", "text-only")),
    "json": (to_json, Capability("json", "JSON (text + geometry + tables)", "application/json", "json", "structural")),
    "csv": (to_csv, Capability("csv", "CSV (tables)", "text/csv", "csv", "structural")),
    "xlsx": (to_xlsx, Capability("xlsx", "Excel (tables)", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "xlsx", "structural")),
    "docx": (to_docx, Capability("docx", "Word", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "docx", "text-only")),
    "png": (lambda d, **k: to_images(d, "png", **k), Capability("png", "PNG images (zip)", "application/zip", "zip", "raster")),
    "jpg": (lambda d, **k: to_images(d, "jpg", **k), Capability("jpg", "JPEG images (zip)", "application/zip", "zip", "raster")),
    "webp": (lambda d, **k: to_images(d, "webp", **k), Capability("webp", "WebP images (zip)", "application/zip", "zip", "raster")),
    "pptx": (to_pptx, Capability(
        "pptx", "PowerPoint (one slide per page)",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "pptx", "raster")),
}

# Every from-PDF target is now produced in-process. LibreOffice is still
# required in the other direction, for turning Word and Excel files into PDFs.
_OFFICE_TARGETS: tuple = ()


def capabilities() -> List[Capability]:
    result = [capability for _, capability in _FROM_PDF.values()]
    office_ready = _office_binary() is not None

    for target in _OFFICE_TARGETS:
        result.append(Capability(
            target, "PowerPoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx", "structural",
            available=office_ready,
            reason=None if office_ready else OFFICE_UNAVAILABLE,
        ))
    return result


def convert(data: bytes, target: str, *, filename: str = "document",
            scale: float = 2.0) -> ConversionResult:
    """Convert a PDF to the requested format."""
    key = target.lower().strip()
    if key in ("jpeg",):
        key = "jpg"
    if key in ("md",):
        key = "markdown"

    if key not in _FROM_PDF:
        available = ", ".join(sorted(_FROM_PDF))
        raise PDFEngineError(
            f"Cannot convert to '{target}'. Available targets: {available}."
        )

    converter, capability = _FROM_PDF[key]

    try:
        payload, count, warnings = (
            converter(data, scale=scale) if capability.fidelity == "raster"
            else converter(data)
        )
    except (PDFEngineError, PasswordRequired):
        raise
    except Exception as exc:
        raise PDFEngineError(f"The conversion failed: {exc}") from exc

    if not payload:
        raise PDFEngineError("The conversion produced an empty file.")

    stem = filename.rsplit(".", 1)[0] or "document"
    return ConversionResult(
        data=payload,
        filename=f"{stem}.{capability.extension}",
        media_type=capability.media_type,
        target=key,
        fidelity=capability.fidelity,
        note=FIDELITY[capability.fidelity],
        pages=count,
        warnings=warnings,
    )


TO_PDF_EXTENSIONS = {
    "txt": "text", "md": "text", "markdown": "text", "text": "text",
    "csv": "csv",
    "png": "image", "jpg": "image", "jpeg": "image",
    "tif": "image", "tiff": "image", "webp": "image", "bmp": "image", "gif": "image",
    "docx": "office", "doc": "office", "xlsx": "office", "xls": "office",
    "pptx": "office", "ppt": "office", "odt": "office", "rtf": "office",
    "html": "office", "htm": "office",
}


def to_pdf(raw: bytes, source_extension: str, *, title: str = "") -> bytes:
    """Create a PDF from a supported source file."""
    extension = source_extension.lower().lstrip(".")
    kind = TO_PDF_EXTENSIONS.get(extension)

    if kind is None:
        raise PDFEngineError(
            f"Cannot create a PDF from a .{extension} file. Supported: "
            f"{', '.join(sorted(TO_PDF_EXTENSIONS))}."
        )

    if kind == "text":
        return text_to_pdf(raw.decode("utf-8", errors="replace"), title=title)
    if kind == "csv":
        return csv_to_pdf(raw, title=title)
    if kind == "image":
        return images_to_pdf([raw])
    return office_to_pdf(raw, extension)
