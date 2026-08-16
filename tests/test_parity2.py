"""
The second parity batch: links, attachments, Bates, scan cleanup, PowerPoint.

Two of these are as much security features as conveniences. A link can carry a
javascript: action and an attachment can carry an executable, and the scanner
already reports both — so the tools that add them refuse to be the way such a
thing gets in.
"""
import io
import random

import pytest
from PIL import Image, ImageDraw
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

from docintel.pdf import attachments, convert, enhance, links
from docintel.pdf.engine import PDFEngineError


def make_pdf(pages: int = 2) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    for number in range(1, pages + 1):
        pdf.setFont("Helvetica", 14)
        pdf.drawString(72, 700, f"Page {number} of the report")
        pdf.showPage()
    pdf.save()
    return buffer.getvalue()


def scan_font():
    """A font with realistic stroke weight.

    PIL's built-in font draws one-pixel strokes, which are genuinely
    indistinguishable from dirt — no despeckle can keep one and drop the
    other. Text on a real 200 dpi scan is three or four pixels thick, so the
    fixture uses a proper face and falls back only if none is installed.
    """
    from PIL import ImageFont
    for name in ("arial.ttf", "DejaVuSans.ttf", "LiberationSans-Regular.ttf"):
        try:
            return ImageFont.truetype(name, 34)
        except OSError:
            continue
    return ImageFont.load_default()


def make_scan(*, skew: float = -4.0, specks: int = 3000) -> bytes:
    image = Image.new("RGB", (1275, 1650), "white")
    draw = ImageDraw.Draw(image)
    font = scan_font()
    for index, line in enumerate(
            ["QUARTERLY REPORT", "Revenue increased by 12%", "Costs held flat"]):
        draw.text((150, 260 + index * 80), line, fill="black", font=font)
    image = image.rotate(skew, resample=Image.BICUBIC, fillcolor="white")

    random.seed(1)
    dirt = ImageDraw.Draw(image)
    for _ in range(specks):
        dirt.point((random.randrange(1275), random.randrange(1650)), fill="black")

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    pdf.drawImage(ImageReader(image), 0, 0, width=letter[0], height=letter[1])
    pdf.save()
    return buffer.getvalue()


# --- links -----------------------------------------------------------------

def test_a_link_can_be_added_and_read_back():
    out = links.add_link(make_pdf(), page=1,
                         rect={"x": 72, "y": 86, "width": 90, "height": 14},
                         url="https://example.com/report")

    found = links.list_links(out)
    assert len(found) == 1
    assert found[0]["kind"] == "uri"
    assert found[0]["target"] == "https://example.com/report"
    assert found[0]["page"] == 1


def test_javascript_links_are_refused():
    """A link is a place a document can hide behaviour."""
    for url in ("javascript:alert(1)", "file:///etc/passwd", "data:text/html,x"):
        with pytest.raises(PDFEngineError, match="not allowed"):
            links.add_link(make_pdf(), page=1,
                           rect={"x": 1, "y": 1, "width": 10, "height": 10},
                           url=url)


def test_a_link_needs_an_area():
    with pytest.raises(PDFEngineError, match="width and a height"):
        links.add_link(make_pdf(), page=1,
                       rect={"x": 1, "y": 1, "width": 0, "height": 0},
                       url="https://example.com")


def test_links_can_be_removed():
    withlink = links.add_link(make_pdf(), page=1,
                              rect={"x": 72, "y": 86, "width": 90, "height": 14},
                              url="https://example.com")
    out, removed = links.remove_links(withlink)

    assert removed == 1
    assert links.list_links(out) == []


def test_removing_links_from_a_document_with_none_says_so():
    with pytest.raises(PDFEngineError, match="no links"):
        links.remove_links(make_pdf())


# --- attachments -----------------------------------------------------------

def test_a_file_can_be_attached_listed_and_extracted():
    out = attachments.attach(make_pdf(), "figures.csv", b"a,b\n1,2\n", "Backing data")

    listed = attachments.list_attachments(out)
    assert [a["name"] for a in listed] == ["figures.csv"]

    payload, name = attachments.extract(out, "figures.csv")
    assert payload == b"a,b\n1,2\n"
    assert name == "figures.csv"


def test_executables_are_refused():
    """Documents carrying executables are treated as malicious, correctly."""
    for name in ("payload.exe", "run.bat", "thing.js", "installer.msi"):
        with pytest.raises(PDFEngineError, match="will not be embedded"):
            attachments.attach(make_pdf(), name, b"MZ")


def test_a_path_is_stripped_from_an_attachment_name():
    """../../etc/passwd must never travel as a path."""
    assert attachments.safe_name("../../etc/passwd") == "passwd"
    assert attachments.safe_name(r"C:\\Windows\\system32\\evil.txt") == "evil.txt"

    out = attachments.attach(make_pdf(), "../../secret.txt", b"x")
    assert [a["name"] for a in attachments.list_attachments(out)] == ["secret.txt"]


def test_attaching_the_same_name_twice_is_refused():
    once = attachments.attach(make_pdf(), "notes.txt", b"one")
    with pytest.raises(PDFEngineError, match="already attached"):
        attachments.attach(once, "notes.txt", b"two")


def test_an_attachment_can_be_removed():
    once = attachments.attach(make_pdf(), "notes.txt", b"one")
    out, removed = attachments.remove(once, "notes.txt")

    assert removed == 1
    assert attachments.list_attachments(out) == []


def test_extracting_something_that_is_not_there():
    with pytest.raises(PDFEngineError, match="not attached"):
        attachments.extract(make_pdf(), "nothing.txt")


# --- scan enhancement ------------------------------------------------------

def test_deskew_straightens_a_tilted_scan():
    out, report = enhance.enhance(make_scan(skew=-4.0), despeckle=False,
                                  contrast=False)

    corrected = report[0]["skew_corrected"]
    assert corrected > 2.0, f"barely corrected: {corrected}"
    assert corrected < 6.0, f"over-rotated: {corrected}"


def test_a_straight_page_is_left_alone():
    """Correcting a page that is already level would only blur it."""
    _, report = enhance.enhance(make_scan(skew=0.0, specks=0), despeckle=False,
                                contrast=False)

    assert abs(report[0]["skew_corrected"]) < 0.5


def test_despeckle_removes_dirt_without_eating_the_text():
    """A filter wide enough to clear dirt must not erode the glyphs."""
    import numpy as np
    from docintel.pdf.render import render_page

    scan = make_scan(skew=0.0, specks=3000)
    out, _ = enhance.enhance(scan, deskew=False, contrast=False)

    def measure(data):
        page = np.asarray(Image.open(io.BytesIO(
            render_page(data, 1, scale=2.0, fmt="png").data)).convert("L"))
        return ((page[1200:2100, :] < 200).mean(),      # dirt, below the text
                (page[380:760, 150:900] < 128).mean())  # ink, in the text band

    dirt_before, ink_before = measure(scan)
    dirt_after, ink_after = measure(out)

    assert dirt_after < dirt_before * 0.2, "the dirt survived"
    assert ink_after > ink_before * 0.3, "the text was eaten"


def test_enhancement_must_be_asked_for():
    with pytest.raises(PDFEngineError, match="No enhancement"):
        enhance.enhance(make_scan(), deskew=False, despeckle=False,
                        contrast=False, binarise=False)


def test_a_page_that_does_not_exist_is_refused():
    with pytest.raises(PDFEngineError, match="do not exist"):
        enhance.enhance(make_scan(), pages=[9])


# --- powerpoint ------------------------------------------------------------

def test_pptx_has_one_slide_per_page_carrying_the_text():
    """Opened with the library that owns the format; bytes prove nothing."""
    from pptx import Presentation

    result = convert.convert(make_pdf(pages=3), "pptx", filename="deck")
    deck = Presentation(io.BytesIO(result.data))

    assert result.pages == 3
    assert len(deck.slides) == 3
    for index, slide in enumerate(deck.slides, start=1):
        assert any(shape.shape_type == 13 for shape in slide.shapes), "no picture"
        notes = slide.notes_slide.notes_text_frame.text
        assert f"Page {index}" in notes


def test_pptx_declares_itself_raster():
    """It is a picture of each page, and the fidelity must say so."""
    result = convert.convert(make_pdf(pages=1), "pptx")

    assert result.fidelity == "raster"
    assert "no longer selectable" in result.fidelity_note


def test_pptx_is_offered_as_available():
    targets = {c.target: c for c in convert.capabilities()}

    assert "pptx" in targets
    assert targets["pptx"].available is True


# --- through the API -------------------------------------------------------

def test_bates_numbering_stamps_a_padded_sequence(alice):
    document = alice.upload(make_pdf(pages=3), "case.pdf").json()["document"]["id"]

    response = alice.post(f"/api/v1/documents/{document}/bates",
                          json={"prefix": "ACME-", "digits": 6, "start_at": 1})

    assert response.status_code == 200, response.text
    assert "ACME-000001" in response.json()["note"]

    from docintel.pdf.text import extract
    latest = alice.get(f"/api/v1/documents/{document}/download").content
    text = " ".join(p.text for p in extract(latest))
    assert "ACME-000001" in text
    assert "ACME-000003" in text


def test_enhance_refuses_a_document_that_has_real_text(alice):
    """Rasterising a text document destroys the thing that makes it useful."""
    document = alice.upload(make_pdf(), "text.pdf").json()["document"]["id"]

    response = alice.post(f"/api/v1/documents/{document}/enhance", json={})

    assert response.status_code == 400
    assert "already has a real text layer" in response.json()["detail"]


def test_enhance_proceeds_when_confirmed(alice):
    document = alice.upload(make_pdf(), "text.pdf").json()["document"]["id"]

    response = alice.post(f"/api/v1/documents/{document}/enhance",
                          json={"confirm_rasterise": True, "deskew": False})

    assert response.status_code == 200, response.text
    assert "run OCR" in response.json()["note"]


def test_link_endpoints_round_trip(alice):
    document = alice.upload(make_pdf(), "doc.pdf").json()["document"]["id"]

    added = alice.post(f"/api/v1/documents/{document}/links", json={
        "page": 1, "rect": {"x": 72, "y": 86, "width": 90, "height": 14},
        "url": "https://example.com"})
    assert added.status_code == 200, added.text

    listed = alice.get(f"/api/v1/documents/{document}/links").json()
    assert listed["count"] == 1

    removed = alice.post(f"/api/v1/documents/{document}/links/remove", json={})
    assert removed.status_code == 200
    assert alice.get(f"/api/v1/documents/{document}/links").json()["count"] == 0


def test_attachment_endpoints_round_trip(alice):
    document = alice.upload(make_pdf(), "doc.pdf").json()["document"]["id"]

    added = alice.client.post(
        f"/api/v1/documents/{document}/attachments",
        headers=alice.headers,
        files={"file": ("data.csv", b"a,b\n1,2\n", "text/csv")},
    )
    assert added.status_code == 201, added.text

    listed = alice.get(f"/api/v1/documents/{document}/attachments").json()
    assert [a["name"] for a in listed["attachments"]] == ["data.csv"]

    fetched = alice.get(f"/api/v1/documents/{document}/attachments/data.csv")
    assert fetched.status_code == 200
    assert fetched.content == b"a,b\n1,2\n"


def test_the_api_refuses_to_attach_an_executable(alice):
    document = alice.upload(make_pdf(), "doc.pdf").json()["document"]["id"]

    response = alice.client.post(
        f"/api/v1/documents/{document}/attachments",
        headers=alice.headers,
        files={"file": ("payload.exe", b"MZ\x90\x00", "application/octet-stream")},
    )

    assert response.status_code == 400
    assert "will not be embedded" in response.json()["detail"]
