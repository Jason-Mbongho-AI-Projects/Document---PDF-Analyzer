"""Reading office formats when LibreOffice is not installed.

The fallback is not a layout engine: it recovers the words and the tables so
a document can still be read, searched and asked questions about on a server
that has no LibreOffice. These tests pin that promise, and the limits of it.
"""
import io

import pypdfium2 as pdfium
import pytest

import pdf_corpus as corpus
from docintel.pdf import convert
from docintel.pdf.engine import PDFEngineError


def text_of(data: bytes) -> str:
    document = pdfium.PdfDocument(data)
    return "\n".join(page.get_textpage().get_text_range() for page in document)


def without_libreoffice(monkeypatch):
    """Make the real converter unavailable, as it is on a bare server."""
    monkeypatch.setattr(convert, "_office_binary", lambda: None)


def test_a_word_document_keeps_its_words(monkeypatch):
    without_libreoffice(monkeypatch)
    text = text_of(convert.to_pdf(corpus.small_docx(), "docx"))
    assert corpus.DOCX_SENTENCE in text
    assert "Quarterly review" in text


def test_a_workbook_keeps_its_cells(monkeypatch):
    without_libreoffice(monkeypatch)
    text = text_of(convert.to_pdf(corpus.small_xlsx(), "xlsx"))
    assert "Region" in text and "128400" in text


def test_a_deck_keeps_one_block_per_slide(monkeypatch):
    without_libreoffice(monkeypatch)
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for title in ("Roadmap", "Budget"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        box.text_frame.paragraphs[0].add_run().text = f"Notes for {title.lower()}."
    buffer = io.BytesIO()
    presentation.save(buffer)

    text = text_of(convert.to_pdf(buffer.getvalue(), "pptx"))
    assert "Slide 1" in text and "Slide 2" in text
    assert "Roadmap" in text and "Notes for budget." in text


def test_html_is_stripped_to_text(monkeypatch):
    without_libreoffice(monkeypatch)
    source = (b"<html><body><h1>Notice</h1>"
              b"<p>Body &amp; text.</p><ul><li>One</li></ul></body></html>")
    text = text_of(convert.to_pdf(source, "html"))
    assert "Notice" in text
    assert "Body & text." in text
    assert "<p>" not in text


def test_script_content_is_dropped_not_shown(monkeypatch):
    """Script bodies are not document text and must not appear as if they were."""
    without_libreoffice(monkeypatch)
    source = (b"<html><body><p>Visible.</p>"
              b"<script>var secret = 'do-not-print';</script></body></html>")
    text = text_of(convert.to_pdf(source, "html"))
    assert "Visible." in text
    assert "do-not-print" not in text


def test_rtf_control_words_are_removed_and_escapes_decoded(monkeypatch):
    without_libreoffice(monkeypatch)
    source = rb"{\rtf1\ansi\deff0 Caf\'e9 opens Monday.\par Second line.}"
    text = text_of(convert.to_pdf(source, "rtf"))
    assert "Caf\u00e9 opens Monday." in text
    assert "Second line." in text
    assert "rtf1" not in text and "ansi" not in text


def test_the_result_says_the_layout_was_not_preserved(monkeypatch):
    """Someone comparing this to the original deserves to know why it differs."""
    without_libreoffice(monkeypatch)
    text = text_of(convert.to_pdf(corpus.small_docx(), "docx"))
    assert "LibreOffice" in text


def test_an_empty_document_still_produces_a_readable_pdf(monkeypatch):
    without_libreoffice(monkeypatch)
    from docx import Document as DocxDocument

    buffer = io.BytesIO()
    DocxDocument().save(buffer)
    data = convert.to_pdf(buffer.getvalue(), "docx")
    assert data.startswith(b"%PDF-")
    assert "no extractable text" in text_of(data)


@pytest.mark.parametrize("extension", ["doc", "xls", "ppt", "odt"])
def test_legacy_formats_have_no_fallback_and_say_so(monkeypatch, extension):
    """We cannot read these ourselves; refusing beats inventing content."""
    without_libreoffice(monkeypatch)
    with pytest.raises(PDFEngineError) as raised:
        convert.to_pdf(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 64, extension)
    assert "LibreOffice" in str(raised.value)


def test_a_corrupt_docx_is_reported_not_crashed(monkeypatch):
    without_libreoffice(monkeypatch)
    with pytest.raises(PDFEngineError):
        convert.to_pdf(b"PK\x03\x04" + b"\x00" * 100, "docx")


def test_libreoffice_is_still_preferred_when_present():
    """The fallback must not take over from the converter that does layout."""
    if convert._office_binary() is None:
        pytest.skip("LibreOffice is not installed on this machine")
    text = text_of(convert.to_pdf(corpus.small_docx(), "docx"))
    assert corpus.DOCX_SENTENCE in text
    assert "LibreOffice" not in text
